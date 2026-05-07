"""
Agrega 5 filminas al PPTX de Unidad 4 - Data Science.pptx
Temas: MCAR/MAR/MNAR | loc/iloc/Máscaras | Resampling/Rolling | sklearn Imputers | Plotly Express
"""
import copy
from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SRC  = r"c:\Users\Turi\Desktop\Data Science I\Data Science 1\Clase04\Unidad 4 - Data Science.pptx"
DEST = r"c:\Users\Turi\Desktop\Data Science I\Data Science 1\Clase04\Unidad 4 - Data Science.pptx"

prs = Presentation(SRC)

# ─── Layout de referencia (mismo que las slides de contenido) ────────────────
LAYOUT_IDX = 0   # SECTION_HEADER_1_1_1_1_1_1_1_1

# ─── Dimensiones del slide (EMU) ─────────────────────────────────────────────
SW = prs.slide_width   # 9144000
SH = prs.slide_height  # 5143500

# ─── Constantes de posición ───────────────────────────────────────────────────
T_LEFT  = Emu(473350)   # 49 px
T_TOP   = Emu(619525)   # 65 px
T_W     = Emu(8141100)  # 854 px
T_H     = Emu(1600000)  # 168 px

BODY_TOP  = Emu(2000000)
BODY_H    = Emu(2950000)
COL1_W    = Emu(3900000)
COL2_LEFT = Emu(4550000)
COL2_W    = Emu(3900000)
FULL_W    = Emu(8141100)


# ─── Helpers ─────────────────────────────────────────────────────────────────
def add_slide():
    layout = prs.slide_layouts[LAYOUT_IDX]
    return prs.slides.add_slide(layout)


def add_textbox(slide, left, top, width, height, text, size=14,
                bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)
    return txb


def add_para(tf, text, size=13, bold=False, color=None,
             bullet=True, indent=0, align=PP_ALIGN.LEFT):
    """Agrega un párrafo a un text frame existente."""
    from pptx.oxml.ns import qn
    p = tf.add_paragraph()
    p.alignment = align
    if bullet and text.strip():
        p.level = indent
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)
    return p


def add_code_box(slide, left, top, width, height, lines):
    """Caja de código con fondo gris oscuro."""
    from pptx.util import Emu as E
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    # fondo gris oscuro
    fill = txb.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x2B, 0x2B, 0x2B)
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(11)
        run.font.bold = False
        run.font.color.rgb = RGBColor(0xF8, 0xF8, 0xF2)  # blanco suave
        # fuente monoespaciada
        from pptx.oxml.ns import qn
        from lxml import etree
        rPr = run._r.get_or_add_rPr()
        latin = etree.SubElement(rPr, qn('a:latin'))
        latin.set('typeface', 'Courier New')
    return txb


def title_color():
    return (0x00, 0x00, 0x00)   # negro (igual que slides originales)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – MCAR / MAR / MNAR
# ═══════════════════════════════════════════════════════════════════════════════
s1 = add_slide()

# Título
add_textbox(s1, T_LEFT, T_TOP, T_W, T_H,
            "Tipos de valores ausentes: MCAR, MAR y MNAR",
            size=28, bold=True, color=title_color())

# Subtítulo / intro
add_textbox(s1, T_LEFT, Emu(1900000), T_W, Emu(600000),
            "No todos los datos faltantes son iguales. El tipo de ausencia determina la estrategia de imputación.",
            size=13, color=(0x33, 0x33, 0x33))

# Tabla de 3 tipos (3 cajas lado a lado)
box_w = Emu(2600000)
gap   = Emu(170000)
tops  = Emu(2550000)
bh    = Emu(2300000)
types = [
    ("MCAR", "Missing Completely\nAt Random",
     "Los nulos son aleatorios:\nno dependen de ninguna\nvariable del dataset.\n\nEjemplo: sensor que falla\nen momentos aleatorios.\n\nAcción: imputar o\neliminar sin sesgo.",
     (0x1A, 0x73, 0xC8)),
    ("MAR", "Missing At Random",
     "Los nulos dependen de\notras variables observadas.\n\nEjemplo: personas mayores\nno declaran su ingreso.\n\nAcción: imputar\ncondicionado al grupo.",
     (0xE6, 0x7E, 0x22)),
    ("MNAR", "Missing Not\nAt Random",
     "Los nulos dependen de\nla variable misma.\n\nEjemplo: pacientes graves\nno responden encuestas.\n\nAcción: modelos de\nselección o recolección.",
     (0xC0, 0x39, 0x2B)),
]

for i, (tag, subtitle, body, col) in enumerate(types):
    lft = T_LEFT + Emu(i * (box_w.emu + gap.emu))
    # Fondo de color de cada caja
    shape = s1.shapes.add_textbox(lft, tops, box_w, bh)
    shape.word_wrap = True
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*col)
    tf = shape.text_frame
    tf.word_wrap = True

    # Línea 1: sigla grande
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = tag
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Línea 2: nombre completo
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.size = Pt(11)
    r2.font.bold = False
    r2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Separador
    p3 = tf.add_paragraph()
    p3.add_run().text = ""

    # Cuerpo
    for line in body.split("\n"):
        px = tf.add_paragraph()
        px.alignment = PP_ALIGN.LEFT
        rx = px.add_run()
        rx.text = line
        rx.font.size = Pt(11)
        rx.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – loc / iloc / Máscaras Booleanas
# ═══════════════════════════════════════════════════════════════════════════════
s2 = add_slide()

add_textbox(s2, T_LEFT, T_TOP, T_W, T_H,
            "Selección de datos: loc, iloc y máscaras booleanas",
            size=28, bold=True, color=title_color())

# Columna izquierda: tabla loc vs iloc
add_textbox(s2, T_LEFT, Emu(1850000), COL1_W, Emu(500000),
            "loc  vs  iloc", size=15, bold=True, color=(0x1A, 0x73, 0xC8))

table_lines = [
    "           loc               iloc",
    "─────────────────────────────────────",
    "Referencia  Etiquetas         Posición",
    "Extremos    Ambos inclusive   Final exc.",
    "Uso típico  df.loc[0:5,'col'] df.iloc[0:5,0]",
]
add_code_box(s2, T_LEFT, Emu(2350000), COL1_W, Emu(1400000), table_lines)

# Columna derecha: máscaras
add_textbox(s2, COL2_LEFT, Emu(1850000), COL2_W, Emu(500000),
            "Máscaras booleanas", size=15, bold=True, color=(0xE6, 0x7E, 0x22))

mask_lines = [
    "# AND: ambas condiciones True",
    "mask = (df['region']=='Norte') \\",
    "     & (df['ventas'] > 5000)",
    "",
    "# OR: al menos una True",
    "mask2 = (df['edad'] < 30) \\",
    "      | (df['ventas'] > 8000)",
    "",
    "# NOT: niega la condición",
    "no_centro = df[~(df['region']=='Centro')]",
    "",
    "# Asignación condicional",
    "df.loc[mask, 'segmento'] = 'Premium'",
]
add_code_box(s2, COL2_LEFT, Emu(2350000), COL2_W, Emu(2350000), mask_lines)

# Nota al pie
add_textbox(s2, T_LEFT, Emu(4750000), T_W, Emu(300000),
            "⚠  Usar & | ~ (nunca and/or/not). Cada condición entre paréntesis es obligatorio.",
            size=11, color=(0x88, 0x00, 0x00), bold=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – Resampling y Ventanas Móviles
# ═══════════════════════════════════════════════════════════════════════════════
s3 = add_slide()

add_textbox(s3, T_LEFT, T_TOP, T_W, T_H,
            "Series de tiempo: Resampling y Ventanas Móviles",
            size=28, bold=True, color=title_color())

# Intro
add_textbox(s3, T_LEFT, Emu(1850000), T_W, Emu(450000),
            "Técnicas para cambiar la granularidad temporal y suavizar series con ruido.",
            size=13, color=(0x33, 0x33, 0x33))

# Columna izquierda: resample
add_textbox(s3, T_LEFT, Emu(2400000), COL1_W, Emu(450000),
            "resample()  — cambiar frecuencia", size=13, bold=True, color=(0x1A, 0x73, 0xC8))

res_lines = [
    "# Horas → días (downsample)",
    "ventas_d = serie_h.resample('D').sum()",
    "",
    "# Días → semanas",
    "ventas_w = serie_h.resample('W').mean()",
    "",
    "# Meses → días (upsample)",
    "ventas_up = serie_m.resample('D').ffill()",
]
add_code_box(s3, T_LEFT, Emu(2900000), COL1_W, Emu(1900000), res_lines)

# Columna derecha: rolling / ewm / expanding
add_textbox(s3, COL2_LEFT, Emu(2400000), COL2_W, Emu(450000),
            "Suavizado — rolling / ewm / expanding", size=13, bold=True, color=(0xE6, 0x7E, 0x22))

roll_lines = [
    "# Media móvil simple (ventana fija)",
    "rolling = serie.rolling(window=7).mean()",
    "# → NaN en las primeras 6 filas",
    "",
    "# Exponencial (más peso a lo reciente)",
    "ewma = serie.ewm(span=7).mean()",
    "# → sin NaN desde la primera fila",
    "",
    "# Media acumulada (ventana creciente)",
    "acum = serie.expanding().mean()",
]
add_code_box(s3, COL2_LEFT, Emu(2900000), COL2_W, Emu(1900000), roll_lines)

# Nota
add_textbox(s3, T_LEFT, Emu(4870000), T_W, Emu(250000),
            "Rolling pondera igual todos los valores. EWMA da más peso a los recientes → reacciona más rápido.",
            size=11, color=(0x44, 0x44, 0x44))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – Imputación avanzada con sklearn
# ═══════════════════════════════════════════════════════════════════════════════
s4 = add_slide()

add_textbox(s4, T_LEFT, T_TOP, T_W, T_H,
            "Imputación avanzada con scikit-learn",
            size=28, bold=True, color=title_color())

add_textbox(s4, T_LEFT, Emu(1850000), T_W, Emu(430000),
            "sklearn provee imputadores que se integran con Pipeline para evitar data leakage.",
            size=13, color=(0x33, 0x33, 0x33))

# Columna izquierda: SimpleImputer
add_textbox(s4, T_LEFT, Emu(2380000), COL1_W, Emu(450000),
            "SimpleImputer  — univariante", size=13, bold=True, color=(0x1A, 0x73, 0xC8))

si_lines = [
    "from sklearn.impute import SimpleImputer",
    "",
    "# strategy: 'mean' | 'median'",
    "#           'most_frequent' | 'constant'",
    "imp = SimpleImputer(strategy='median')",
    "X_clean = imp.fit_transform(X)",
    "",
    "# ColumnTransformer: distintas",
    "# estrategias por tipo de columna",
    "from sklearn.compose import ColumnTransformer",
    "ct = ColumnTransformer([",
    "  ('num', SimpleImputer('median'), num_cols),",
    "  ('cat', SimpleImputer('most_frequent'), cat_cols),",
    "])",
]
add_code_box(s4, T_LEFT, Emu(2870000), COL1_W, Emu(2000000), si_lines)

# Columna derecha: IterativeImputer
add_textbox(s4, COL2_LEFT, Emu(2380000), COL2_W, Emu(450000),
            "IterativeImputer  — MICE multivariante", size=13, bold=True, color=(0xE6, 0x7E, 0x22))

ii_lines = [
    "from sklearn.impute import IterativeImputer",
    "# (experimental en sklearn < 1.0)",
    "",
    "# Modela cada columna con NaN",
    "# como función de las demás.",
    "# Itera hasta convergencia.",
    "imp_iter = IterativeImputer(",
    "    max_iter=10,",
    "    random_state=0,",
    ")",
    "X_mice = imp_iter.fit_transform(X)",
    "",
    "# Captura correlaciones entre cols",
    "# → mejor que SimpleImputer",
    "# cuando las vars están correlacionadas",
]
add_code_box(s4, COL2_LEFT, Emu(2870000), COL2_W, Emu(2000000), ii_lines)

add_textbox(s4, T_LEFT, Emu(4920000), T_W, Emu(250000),
            "Siempre usar fit() solo sobre train set. Pipeline lo garantiza automáticamente.",
            size=11, color=(0x88, 0x00, 0x00), bold=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – Plotly Express
# ═══════════════════════════════════════════════════════════════════════════════
s5 = add_slide()

add_textbox(s5, T_LEFT, T_TOP, T_W, T_H,
            "Visualizaciones interactivas con Plotly Express",
            size=28, bold=True, color=title_color())

add_textbox(s5, T_LEFT, Emu(1850000), T_W, Emu(400000),
            "A diferencia de Matplotlib, Plotly genera gráficos HTML con hover, zoom y export nativos.",
            size=13, color=(0x33, 0x33, 0x33))

# Tres cajas de gráficos
chart_w = Emu(2500000)
chart_gap = Emu(220000)
chart_top = Emu(2380000)
chart_h = Emu(2500000)

charts = [
    ("📊  Bar chart", (0x1A, 0x73, 0xC8),
     ["import plotly.express as px",
      "",
      "fig = px.bar(",
      "  df.sort_values('ventas'),",
      "  x='categoria',",
      "  y='ventas',",
      "  color='crecimiento',",
      "  color_continuous_scale=",
      "    'RdYlGn',",
      ")",
      "fig.show()"]),
    ("🍩  Donut chart", (0x27, 0xAE, 0x60),
     ["fig = px.pie(",
      "  df,",
      "  names='categoria',",
      "  values='ventas',",
      "  hole=0.4,  # donut",
      ")",
      "fig.update_traces(",
      "  textinfo=",
      "   'percent+label'",
      ")",
      "fig.show()"]),
    ("🌳  Treemap + Export", (0x8E, 0x44, 0xAD),
     ["fig = px.treemap(",
      "  df,",
      "  path=['sector',",
      "        'categoria'],",
      "  values='ventas',",
      "  color='margen',",
      ")",
      "fig.show()",
      "",
      "# Exportar a HTML",
      "fig.write_html(",
      "  'reporte.html')"]),
]

for i, (label, col, code) in enumerate(charts):
    lft = T_LEFT + Emu(i * (chart_w.emu + chart_gap.emu))

    # Encabezado de la caja
    hdr = s5.shapes.add_textbox(lft, chart_top, chart_w, Emu(550000))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = RGBColor(*col)
    tf_h = hdr.text_frame
    p = tf_h.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Código
    add_code_box(s5, lft, chart_top + Emu(580000), chart_w,
                 chart_h - Emu(580000), code)

add_textbox(s5, T_LEFT, Emu(4930000), T_W, Emu(250000),
            "fig.write_html('archivo.html') guarda el gráfico completo sin necesitar servidor ni Python.",
            size=11, color=(0x44, 0x44, 0x44))


# ─── Insertar las 5 slides antes de "¿Preguntas?" (slide 47, índice 46) ──────
# Las slides nuevas se agregan al final; las movemos a la posición correcta.
from pptx.oxml.ns import qn

sldIdLst = prs.slides._sldIdLst
all_ids  = list(sldIdLst)

# Índices de las 5 slides nuevas (fueron agregadas al final)
total = len(prs.slides)
new_indices = list(range(total - 5, total))   # últimas 5

# Queremos insertarlas en la posición 46 (antes de "¿Preguntas?")
insert_pos = 46

# Extraer los elementos de las 5 slides nuevas
new_elems = [all_ids[i] for i in new_indices]

# Removerlos de su posición actual
for el in new_elems:
    sldIdLst.remove(el)

# Insertar en la posición deseada
ref = sldIdLst[insert_pos]
for el in reversed(new_elems):
    sldIdLst.insert(list(sldIdLst).index(ref), el)

prs.save(DEST)
print(f"Guardado: {DEST}")
print(f"Total slides: {len(prs.slides)}")
print(f"Slides nuevas insertadas en posición {insert_pos + 1}")
